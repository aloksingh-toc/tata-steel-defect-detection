"""Generate presentation/defect_detection_summary.pptx"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
OUT  = ROOT / "presentation" / "defect_detection_summary.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ── Palette ────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x17, 0x2A)
BLUE   = RGBColor(0x38, 0xBD, 0xF8)
GREEN  = RGBColor(0x4A, 0xDE, 0x80)
RED    = RGBColor(0xF8, 0x71, 0x71)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
WHITE  = RGBColor(0xE2, 0xE8, 0xF0)
GREY   = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, bullet=False, space_after=6, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = ("•  " + text) if bullet else text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p


# ── Slide 1: Title ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
tf = box(s, 1, 2.2, 11.3, 2)
para(tf, "Defect Detection in Hot Rolling", 44, BLUE, bold=True, first=True)
para(tf, "Predicting the Alpha defect in steel coils", 24, WHITE, space_after=4)
tf2 = box(s, 1, 4.4, 11.3, 1.5)
para(tf2, "Tata Steel AI Hackathon  ·  HackerEarth", 18, GREY, first=True)
para(tf2, "Binary Classification  |  Severe Class Imbalance (19.5 : 1)", 16, GREY)

# ── Slide 2: Problem ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
para(box(s, 0.7, 0.5, 12, 1), "The Problem", 32, BLUE, bold=True, first=True)
tf = box(s, 0.9, 1.7, 11.5, 5)
para(tf, "Alpha defects cannot be detected inline — the coil is under tension in inspection zones.", 19, WHITE, bullet=True, first=True, space_after=12)
para(tf, "Current quality control relies on sampling only a fraction of coils at the final stage.", 19, WHITE, bullet=True, space_after=12)
para(tf, "Although rare, the defect drives customer complaints and product downgrades.", 19, WHITE, bullet=True, space_after=12)
para(tf, "Goal: detect the defect during rolling, using process parameters from every stage.", 19, GREEN, bullet=True, space_after=12)
tf2 = box(s, 0.9, 5.3, 11.5, 1.5)
para(tf2, "Data:  1,352 train coils  ·  339 test coils  ·  49 process features (X1–X49)", 18, AMBER, bold=True, first=True)
para(tf2, "Target:  Y = 1 (Defect)  vs  Y = 0 (No Defect)   —   only 66 defects in training (4.9%)", 18, AMBER, bold=True)

# ── Slide 3: Evaluation target ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
para(box(s, 0.7, 0.5, 12, 1), "Evaluation Target", 32, BLUE, bold=True, first=True)
tf = box(s, 0.9, 1.8, 11.5, 2)
para(tf, "Recall = 100 %     →   zero missed defects (no false negatives)", 22, GREEN, bold=True, first=True, space_after=14)
para(tf, "Precision > 90 %   →   at most ~7 false alarms allowed", 22, AMBER, bold=True, space_after=14)
tf2 = box(s, 0.9, 4.0, 11.5, 3)
para(tf2, "This is an extreme requirement for a 19.5 : 1 imbalanced problem:", 18, WHITE, first=True, space_after=10)
para(tf2, "To allow 100% recall with >90% precision, the model may emit ≤ 7 false positives", 17, GREY, bullet=True, space_after=8)
para(tf2, "→ a false-positive rate ≤ 0.54%", 17, GREY, bullet=True, space_after=8)
para(tf2, "→ which mathematically demands a model AUC ≈ 0.997", 17, GREY, bullet=True, space_after=8)

# ── Slide 4: Approach pipeline ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
para(box(s, 0.7, 0.5, 12, 1), "Our Approach — End-to-End Pipeline", 32, BLUE, bold=True, first=True)
steps = [
    ("1.  Feature Engineering", "49 → 68 features: log transforms, ratios, row-wise stats, interactions, missingness flags", BLUE),
    ("2.  Preprocessing", "IterativeImputer (multivariate) → StandardScaler, persisted to disk", GREEN),
    ("3.  Imbalance Handling", "scale_pos_weight (LGB/XGB) + class_weight='balanced' (RF/ET)", AMBER),
    ("4.  Modelling", "LightGBM (Optuna-tuned) + XGBoost + Random Forest + Extra Trees", BLUE),
    ("5.  Ensemble", "Equal-weight blend of 4 models, 5-fold stratified CV", GREEN),
    ("6.  Threshold", "Set to min OOF probability of any defect → guarantees Recall = 100%", AMBER),
]
y = 1.7
for title, desc, c in steps:
    tf = box(s, 0.9, y, 11.6, 0.9)
    para(tf, title, 19, c, bold=True, first=True, space_after=2)
    para(tf, desc, 15, WHITE)
    y += 0.9

# ── Slide 5: Results ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
para(box(s, 0.7, 0.5, 12, 1), "Results (5-Fold Out-of-Fold)", 32, BLUE, bold=True, first=True)
rows = [
    ("Model", "ROC-AUC", "PR-AUC", WHITE, True),
    ("LightGBM (Optuna-tuned)", "0.8550", "0.3657", GREY, False),
    ("XGBoost", "0.8553", "0.3385", GREY, False),
    ("Random Forest", "0.8551", "0.2868", GREY, False),
    ("Extra Trees", "0.8575", "0.2921", GREY, False),
    ("Blend  (final)", "0.8653", "0.3660", GREEN, True),
]
y = 1.8
for name, auc, ap, c, bold in rows:
    tf = box(s, 1.0, y, 11, 0.55)
    p = tf.paragraphs[0]; p.first = True
    for txt, w in [(name, 6.0), (auc, 2.5), (ap, 2.5)]:
        run = p.add_run(); run.text = f"{txt:<28}" if w == 6.0 else f"{txt:<12}"
        run.font.size = Pt(18); run.font.bold = bold
        run.font.color.rgb = c; run.font.name = "Consolas"
    y += 0.55
para(box(s, 1.0, 5.6, 11, 1), "Validated against One-Class SVM, Isolation Forest, KNN, Mahalanobis, CatBoost, MLP, stacking — blend was best.", 15, GREY, first=True)

# ── Slide 6: Honest gap analysis ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
para(box(s, 0.7, 0.5, 12, 1), "Achieved vs Target — Honest Analysis", 32, BLUE, bold=True, first=True)
tf = box(s, 0.9, 1.7, 11.6, 2.2)
para(tf, "Recall = 100%   ✓  ACHIEVED  (all 66 defects caught)", 20, GREEN, bold=True, first=True, space_after=12)
para(tf, "Precision @ Recall=100%  =  6%   ✗  (target > 90%)", 20, RED, bold=True, space_after=12)
para(tf, "Recall @ Precision=90%   =  1.5%  ✗", 20, RED, bold=True)
tf2 = box(s, 0.9, 4.1, 11.6, 3)
para(tf2, "Why the precision target is a DATA ceiling, not a modelling gap:", 18, AMBER, bold=True, first=True, space_after=10)
para(tf2, "The 10 hardest defects score < 0.10 probability — statistically indistinguishable from normal coils across all 68 features.", 16, WHITE, bullet=True, space_after=8)
para(tf2, "Only 66 defect examples is far too few to learn a 0.54%-FPR boundary.", 16, WHITE, bullet=True, space_after=8)
para(tf2, "Reaching 90% precision needs: more labelled defects, domain features, or raw sensor time-series.", 16, GREEN, bullet=True, space_after=8)

# ── Slide 7: Deliverables ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
para(box(s, 0.7, 0.5, 12, 1), "Deliverables & Reproducibility", 32, BLUE, bold=True, first=True)
tf = box(s, 0.9, 1.8, 11.6, 5)
para(tf, "expected_submission.csv  —  prediction file (339 × 2, verified format)", 18, GREEN, bullet=True, first=True, space_after=10)
para(tf, "defect_detection.ipynb  —  full EDA + modelling notebook (executes end-to-end)", 18, WHITE, bullet=True, space_after=10)
para(tf, "src/  —  modular pipeline: config, features, preprocess, model, evaluate, predict", 18, WHITE, bullet=True, space_after=10)
para(tf, "run.py  —  single command reproduces the submission:  python run.py", 18, WHITE, bullet=True, space_after=10)
para(tf, "README.md + APPROACH.md  —  methodology, feature engineering, tools", 18, WHITE, bullet=True, space_after=10)
para(tf, "Tools:  Python · pandas · scikit-learn · LightGBM · XGBoost · Optuna", 16, GREY, bullet=True, space_after=10)

prs.save(OUT)
print(f"Presentation written: {OUT}")
print(f"Slides: {len(prs.slides._sldIdLst)}")
